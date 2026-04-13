"""
AdaptAd — full 10-minute presentation with demo slide and limitations.
Run: python3 make_pptx_v2.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0d, 0x15, 0x1f)
ACCENT  = RGBColor(0x00, 0xd4, 0xff)
GREEN   = RGBColor(0x00, 0xff, 0x88)
ORANGE  = RGBColor(0xff, 0xb8, 0x00)
RED     = RGBColor(0xff, 0x2d, 0x55)
PURPLE  = RGBColor(0xa7, 0x8b, 0xfa)
PINK    = RGBColor(0xf4, 0x72, 0xb6)
WHITE   = RGBColor(0xff, 0xff, 0xff)
SUB     = RGBColor(0x8b, 0x9a, 0xb2)
CARD    = RGBColor(0x14, 0x21, 0x33)
DEEP    = RGBColor(0x08, 0x12, 0x22)
BORDER  = RGBColor(0x20, 0x35, 0x50)
WARN    = RGBColor(0xff, 0xb8, 0x00)

def prs_new():
    p = Presentation()
    p.slide_width  = Inches(13.33)
    p.slide_height = Inches(7.5)
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(s, prs):
    sh = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    sh.fill.solid(); sh.fill.fore_color.rgb = BG
    sh.line.fill.background()

def bar(s, prs, h=Inches(0.055), color=ACCENT):
    sh = s.shapes.add_shape(1, 0, 0, prs.slide_width, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()

def vbar(s, prs, w=Inches(0.055)):
    sh = s.shapes.add_shape(1, 0, 0, w, prs.slide_height)
    sh.fill.solid(); sh.fill.fore_color.rgb = ACCENT
    sh.line.fill.background()

def txt(s, text, l, t, w, h, size=13, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    b = s.shapes.add_textbox(l, t, w, h)
    tf = b.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return b

def chip(s, text, l, t, color=ACCENT, w=None, h=Inches(0.34)):
    cw = w or Inches(max(1.2, len(text) * 0.13))
    sh = s.shapes.add_shape(1, l, t, cw, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    tf = sh.text_frame; p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = BG

def card(s, l, t, w, h, fill=CARD, border=BORDER, bw=0.75):
    sh = s.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border; sh.line.width = Pt(bw)
    return sh

def gene_bar(s, l, t, val, color, max_w=Inches(1.4)):
    w = max(Inches(0.04), max_w * val)
    sh = s.shapes.add_shape(1, l, t, w, Inches(0.22))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()

# ── slide helpers ─────────────────────────────────────────────────────────────

def header(s, section_label, title, prs, subtitle=None):
    bar(s, prs); chip(s, section_label, Inches(0.5), Inches(0.15))
    txt(s, title, Inches(0.5), Inches(0.65), Inches(12.3), Inches(0.8),
        size=30, bold=True)
    if subtitle:
        txt(s, subtitle, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.4),
            size=13, color=SUB)

# ─────────────────────────────────────────────────────────────────────────────
# Slide 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
def s01_title(prs):
    s = blank(prs); bg(s, prs); bar(s, prs, Inches(0.06)); vbar(s, prs)

    txt(s, "AdaptAd", Inches(0.55), Inches(1.3), Inches(10), Inches(1.3),
        size=60, bold=True, color=ACCENT)
    txt(s, "Evolutionary Multi-Agent Ad Decision System for Streaming",
        Inches(0.55), Inches(2.7), Inches(10), Inches(0.65), size=22, color=WHITE)
    txt(s, "CS6170 AI Capstone  ·  Northeastern University  ·  April 2026",
        Inches(0.55), Inches(3.45), Inches(10), Inches(0.45), size=13, color=SUB)
    txt(s, "Craig Roberts  ·  Arzoo Jiwani  ·  Vishwajeet Hogale",
        Inches(0.55), Inches(3.95), Inches(10), Inches(0.4), size=13, color=SUB)

    # Gene bars right side
    GCOLORS = [ACCENT, GREEN, ORANGE, RED, PURPLE, PINK,
               RGBColor(0x34,0xd3,0x99), RGBColor(0xfb,0x92,0x3c)]
    vals = [0.46, 0.92, 0.35, 0.39, 0.08, 0.80, 0.00, 0.66]
    for i, (c, v) in enumerate(zip(GCOLORS, vals)):
        w = Inches(0.12 + 2.2 * v)
        sh = s.shapes.add_shape(1, Inches(10.8),
                                Inches(1.4 + i * 0.5), w, Inches(0.32))
        sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()

    txt(s, "evolved chromosome · 8 genes", Inches(10.7), Inches(5.6),
        Inches(2.5), Inches(0.35), size=9, color=RGBColor(0x2a,0x40,0x58),
        italic=True)

    # Speaker tag
    card(s, Inches(0.55), Inches(5.3), Inches(2.5), Inches(0.42),
         fill=DEEP, border=BORDER, bw=0.5)
    txt(s, "▶  Speaker A — Craig", Inches(0.68), Inches(5.35),
        Inches(2.3), Inches(0.35), size=10, color=SUB)
    return s


# ─────────────────────────────────────────────────────────────────────────────
# Slide 2 — The Problem
# ─────────────────────────────────────────────────────────────────────────────
def s02_problem(prs):
    s = blank(prs); bg(s, prs)
    header(s, "THE PROBLEM", "Streaming ads are hostile — and getting worse", prs)

    stats = [
        ("46 %",  ACCENT, "of viewers use\nad blockers"),
        ("30 %",  RED,    "engagement drop\nfrom overexposure"),
        ("3.5 h", GREEN,  "average daily\nstreaming time"),
        ("$600B", ORANGE, "global ad market\nwith hostile users"),
    ]
    for i, (num, color, label) in enumerate(stats):
        l = Inches(0.35 + i * 3.2)
        card(s, l, Inches(1.75), Inches(3.0), Inches(1.65))
        txt(s, num, l + Inches(0.15), Inches(1.9), Inches(2.7), Inches(0.85),
            size=36, bold=True, color=color)
        txt(s, label, l + Inches(0.15), Inches(2.7), Inches(2.7),
            Inches(0.55), size=11, color=SUB)

    card(s, Inches(0.35), Inches(3.65), Inches(12.65), Inches(1.2), fill=DEEP)
    txt(s, "The industry asks:  \"Which ad gets the most clicks?\"",
        Inches(0.55), Inches(3.78), Inches(12.2), Inches(0.45),
        size=15, bold=True, color=RED)
    txt(s, "We ask:  \"Should we even show an ad right now — and if so, how?\"",
        Inches(0.55), Inches(4.25), Inches(12.2), Inches(0.45),
        size=15, bold=True, color=GREEN)

    problems = [
        "One-size-fits-all frequency caps ignore viewer state (fatigue, mood, content intensity)",
        "Click-through optimisation creates annoying, ill-timed ads that accelerate viewer churn",
        "No system asks whether the current moment is a good time — only which ad to show",
    ]
    for i, p in enumerate(problems):
        txt(s, f"•  {p}", Inches(0.5), Inches(5.1 + i * 0.42),
            Inches(12.3), Inches(0.4), size=12, color=SUB)

    card(s, Inches(0.35), Inches(5.2), Inches(0.06), Inches(1.7),
         fill=RED, border=RED)
    return s


# ─────────────────────────────────────────────────────────────────────────────
# Slide 3 — Our Answer
# ─────────────────────────────────────────────────────────────────────────────
def s03_answer(prs):
    s = blank(prs); bg(s, prs)
    header(s, "OUR ANSWER", "Four decisions. Made intelligently. Every ad break.", prs)

    decisions = [
        ("SHOW",     GREEN,  "Full ad.\nConditions are\nfavorable — user is\nreceptive, ad is\nrelevant."),
        ("SOFTEN",   ACCENT, "Shorter version.\nModerate fit — reduce\nfriction while\nstill generating\nrevenue."),
        ("DELAY",    ORANGE, "Wait for a better\nmoment. Good ad,\nbad timing — do\nnot penalise the\nadvertiser."),
        ("SUPPRESS", RED,    "Skip entirely.\nProtect the viewer.\nHigh fatigue, low\nrelevance, intense\ncontent moment."),
    ]
    for i, (label, color, desc) in enumerate(decisions):
        l = Inches(0.35 + i * 3.2)
        card(s, l, Inches(1.7), Inches(3.05), Inches(5.35))
        chip(s, label, l + Inches(0.15), Inches(1.85), color=color,
             w=Inches(1.8))
        txt(s, desc, l + Inches(0.15), Inches(2.35), Inches(2.75),
            Inches(4.5), size=12, color=SUB)

    txt(s, "Per-opportunity decisions based on: viewer fatigue · content mood · ad relevance · "
           "session depth · time of day · binge state",
        Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.38),
        size=11, color=RGBColor(0x44,0x5a,0x72), italic=True)
    return s


# ─────────────────────────────────────────────────────────────────────────────
# Slide 4 — Architecture
# ─────────────────────────────────────────────────────────────────────────────
def s04_arch(prs):
    s = blank(prs); bg(s, prs)
    header(s, "ARCHITECTURE", "Four-stage pipeline — from viewer context to decision", prs)

    stages = [
        ("① Ingest", ACCENT,
         "UserProfile\nAdCandidate\nContentItem\nSessionContext"),
        ("② Evolve", GREEN,
         "GA optimises\n8-gene chromosome\nover 30 chromosomes\n× 50 generations"),
        ("③ Score", ORANGE,
         "User Advocate\nscores receptivity\nAdvertiser Advocate\nscores value"),
        ("④ Decide", RED,
         "Negotiator maps\ncombined score →\nSHOW / SOFTEN /\nDELAY / SUPPRESS"),
    ]

    for i, (label, color, body) in enumerate(stages):
        l = Inches(0.35 + i * 3.2)
        card(s, l, Inches(1.55), Inches(3.0), Inches(3.85))
        chip(s, label, l + Inches(0.15), Inches(1.68), color=color,
             w=Inches(2.0))
        txt(s, body, l + Inches(0.15), Inches(2.15), Inches(2.7),
            Inches(3.1), size=12, color=SUB)
        if i < 3:
            sh = s.shapes.add_shape(1, l + Inches(3.0),
                                    Inches(3.2), Inches(0.18), Inches(0.22))
            sh.fill.solid(); sh.fill.fore_color.rgb = SUB
            sh.line.fill.background()

    # Bottom row
    bottom = [
        ("LangGraph", PURPLE,
         "Orchestrates agent flow.\nEvolution + decision\ngraphs via LangGraph."),
        ("LLM Explain", RGBColor(0x34,0xd3,0x99),
         "Groq primary (Llama 3.3).\nGemini fallback.\nTemplate if offline."),
        ("SQLite DB", SUB,
         "Decisions, A/B sessions,\nratings, evolution runs.\naiosqlite async."),
        ("React UI", ACCENT,
         "Live evolution chart,\nsession simulator,\nA/B testing panel."),
    ]
    for i, (label, color, body) in enumerate(bottom):
        l = Inches(0.35 + i * 3.2)
        card(s, l, Inches(5.65), Inches(3.0), Inches(1.55), fill=DEEP)
        chip(s, label, l + Inches(0.12), Inches(5.75), color=color,
             w=Inches(1.7))
        txt(s, body, l + Inches(0.12), Inches(6.18), Inches(2.75),
            Inches(0.95), size=10, color=SUB)
    return s


# ─────────────────────────────────────────────────────────────────────────────
# Slide 5 — 8-Gene Chromosome
# ─────────────────────────────────────────────────────────────────────────────
def s05_chrom(prs):
    s = blank(prs); bg(s, prs)
    header(s, "THE CHROMOSOME", "8 genes · all active · evolved, not hand-tuned", prs,
           subtitle="Each gene ∈ [0,1]. The GA tunes them over 50 generations.")

    GCOLORS = [ACCENT, GREEN, ORANGE, RED, PURPLE, PINK,
               RGBColor(0x34,0xd3,0x99), RGBColor(0xfb,0x92,0x3c)]
    genes = [
        ("fatigue_weight",       0.46, "How aggressively session fatigue suppresses ads",         "High = very cautious with tired viewers"),
        ("relevance_weight",     0.92, "How much ad-interest match gates showing an ad",          "High = only show ads to matching users"),
        ("timing_weight",        0.35, "How much time-of-day alignment matters",                  "High = strictly favour preferred watch times"),
        ("frequency_threshold",  0.39, "Base bar for showing any ad at all (maps to 0.35–0.65)", "High = stricter — harder to cross the SHOW line"),
        ("delay_probability",    0.08, "Width of the DELAY zone below the soften threshold",      "High = prefer delaying over suppressing"),
        ("soften_threshold",     0.80, "Width of the SOFTEN zone below the show threshold",       "High = prefer shorter ads over hard skips"),
        ("category_boost",       0.00, "Advertiser weight on category-user relevance match",      "High = heavily rewards category-aligned ads"),
        ("session_depth_factor", 0.66, "How much penalty grows as ads_shown increases",           "High = increasingly cautious deep in session"),
    ]

    row_h = Inches(0.52)
    for i, (name, val, short, effect) in enumerate(genes):
        top = Inches(1.55) + i * row_h
        color = GCOLORS[i]
        gene_bar(s, Inches(0.4), top + Inches(0.15), val, color)
        txt(s, f"{val:.2f}", Inches(0.4), top, Inches(0.65), Inches(0.45),
            size=10, bold=True, color=color)
        txt(s, name, Inches(1.9), top, Inches(2.8), Inches(0.45),
            size=11, bold=True, color=WHITE)
        txt(s, short, Inches(4.8), top, Inches(4.4), Inches(0.45),
            size=10, color=SUB)
        txt(s, effect, Inches(9.3), top, Inches(3.9), Inches(0.45),
            size=10, color=RGBColor(0x44,0x5a,0x72), italic=True)

    card(s, Inches(0.35), Inches(6.88), Inches(12.6), Inches(0.5),
         fill=DEEP, border=BORDER)
    txt(s, "ad_tolerance from UserProfile is also active: high-tolerance users get higher "
           "satisfaction scores for the same SHOW decision — creates realistic heterogeneity "
           "across the 1,000-user synthetic population.",
        Inches(0.55), Inches(6.95), Inches(12.1), Inches(0.38),
        size=10, color=SUB, italic=True)
    return s


# ─────────────────────────────────────────────────────────────────────────────
# Slide 6 — Agent System
# ─────────────────────────────────────────────────────────────────────────────
def s06_agents(prs):
    s = blank(prs); bg(s, prs)
    header(s, "AGENT SYSTEM", "Two advocates score independently. One negotiator decides.", prs)

    panels = [
        ("User Advocate", GREEN, "55% weight", [
            "Scores viewer receptivity at this exact moment",
            "Inputs: session fatigue, ad relevance, time-of-day",
            "match, content mood, binge state, ads_shown",
            "Genes: fatigue_weight · relevance_weight",
            "        timing_weight · session_depth_factor",
            "Output: UA score ∈ [0, 1]",
            "",
            "Higher score = viewer is receptive right now",
        ]),
        ("Advertiser Advocate", ORANGE, "45% weight", [
            "Scores business value of serving this ad here",
            "Inputs: category match, user engagement, primetime",
            "slot, ad priority, seasonal affinity, demo match",
            "Gene: category_boost",
            "Output: ADV score ∈ [0, 1]",
            "",
            "Higher score = advertiser gets good value here",
            "",
        ]),
        ("Negotiator", ACCENT, "combines both", [
            "combined = 0.55 · UA + 0.45 · ADV",
            "Maps score → decision via gene-tuned thresholds:",
            "  combined ≥ show_thresh   → SHOW",
            "  combined ≥ soften_thresh → SOFTEN",
            "  combined ≥ delay_thresh  → DELAY",
            "  otherwise               → SUPPRESS",
            "",
            "All three thresholds are derived from chromosome genes",
        ]),
    ]

    for i, (title, color, sub, lines) in enumerate(panels):
        l = Inches(0.35 + i * 4.3)
        card(s, l, Inches(1.55), Inches(4.1), Inches(5.65))
        chip(s, title, l + Inches(0.15), Inches(1.7), color=color, w=Inches(2.5))
        txt(s, sub, l + Inches(0.15), Inches(2.12), Inches(3.8), Inches(0.35),
            size=10, color=color, italic=True)
        body = "\n".join(lines)
        txt(s, body, l + Inches(0.15), Inches(2.52), Inches(3.82),
            Inches(4.55), size=11, color=SUB)

    txt(s, "LLM (Groq/Gemini/template) adds natural-language explanation after every decision — "
           "\"Skipped: 3 ads already shown, intense scene at 24 min, ad irrelevant to your interests.\"",
        Inches(0.5), Inches(7.07), Inches(12.3), Inches(0.38),
        size=11, color=RGBColor(0x44,0x5a,0x72), italic=True)
    return s


# ─────────────────────────────────────────────────────────────────────────────
# Slide 7 — Genetic Algorithm
# ─────────────────────────────────────────────────────────────────────────────
def s07_ga(prs):
    s = blank(prs); bg(s, prs)
    header(s, "GENETIC ALGORITHM", "Evolving the chromosome over 50 generations", prs)

    steps = [
        ("Init",      ACCENT,  "30 random\nchromosomes\nUniform [0,1]"),
        ("Evaluate",  GREEN,   "Fitness on 1000\nusers × 5 scenarios\nNumPy vectorised"),
        ("Select",    ORANGE,  "3-way tournament:\npick 3, keep best\nfor parent pairs"),
        ("Crossover", ACCENT,  "Uniform: each gene\nfrom A or B\nat 50/50"),
        ("Mutate",    RED,     "Gaussian Δ per gene\nat 15% rate\nclamped [0,1]"),
        ("Elite",     GREEN,   "Top 20% survive\nunchanged into\nnext generation"),
        ("Restart",   ORANGE,  "Stuck 20 gens?\nFresh random pop.\nPreserve best."),
    ]

    bw = Inches(1.78)
    for i, (label, color, desc) in enumerate(steps):
        l = Inches(0.3 + i * 1.87)
        card(s, l, Inches(1.6), bw, Inches(3.2))
        chip(s, label, l + Inches(0.1), Inches(1.72), color=color, w=Inches(1.58))
        txt(s, desc, l + Inches(0.1), Inches(2.15), Inches(1.6),
            Inches(2.55), size=11, color=SUB, align=PP_ALIGN.CENTER)
        if i < 6:
            sh = s.shapes.add_shape(1, l + bw, Inches(2.85),
                                    Inches(0.1), Inches(0.22))
            sh.fill.solid(); sh.fill.fore_color.rgb = SUB
            sh.line.fill.background()

    card(s, Inches(0.3), Inches(5.05), Inches(12.75), Inches(0.85),
         fill=DEEP)
    txt(s, "Fitness  =  0.60 × mean_satisfaction  +  0.40 × mean_revenue",
        Inches(0.5), Inches(5.12), Inches(6.5), Inches(0.38),
        size=14, bold=True, color=ACCENT)
    txt(s, "(60/40 user-revenue split is a deliberate value judgment — see Limitations slide)",
        Inches(7.0), Inches(5.12), Inches(5.8), Inches(0.38),
        size=11, color=ORANGE, italic=True)
    txt(s, "satisfaction measures viewer experience  ·  revenue measures advertiser value  ·  "
           "force-suppress when fatigue > 0.85  ·  ~8 min per full 30×50 run  ·  "
           "convergence window = 15 gen, threshold = 0.0005",
        Inches(0.5), Inches(5.56), Inches(12.3), Inches(0.3),
        size=10, color=SUB)

    card(s, Inches(0.3), Inches(6.05), Inches(12.75), Inches(1.3),
         fill=DEEP)
    txt(s, "Ablation conditions tested alongside full system:",
        Inches(0.5), Inches(6.12), Inches(5.0), Inches(0.35),
        size=12, bold=True, color=WHITE)
    ablations = [
        "Full system (GA + both agents)",
        "GA only — no agents, just fitness",
        "Agents only — no GA, default chromosome",
        "User Advocate only",
        "Advertiser Advocate only",
    ]
    for i, a in enumerate(ablations):
        l = Inches(0.5) if i < 3 else Inches(6.9)
        t = Inches(6.5 + (i % 3) * 0.27)
        txt(s, f"• {a}", l, t, Inches(5.9), Inches(0.26), size=11, color=SUB)
    return s


# ─────────────────────────────────────────────────────────────────────────────
# Slide 8 — Hypotheses
# ─────────────────────────────────────────────────────────────────────────────
def s08_hypo(prs):
    s = blank(prs); bg(s, prs)
    header(s, "HYPOTHESES", "Three testable predictions · Wilcoxon signed-rank · Holm–Bonferroni", prs)

    hyps = [
        ("H1", ACCENT, "Fitness > 0.58",
         "Evolved policy achieves mean composite fitness above 0.58 AND "
         "significantly outperforms all three baselines.\n\n"
         "Baselines:\n"
         "  always_show  —  fitness 0.512\n"
         "  random       —  fitness 0.473\n"
         "  freq_cap_3   —  fitness 0.462\n\n"
         "Test: one-sample Wilcoxon vs 0.58 + paired\n"
         "Wilcoxon vs each baseline with Holm–Bonferroni"),
        ("H2", GREEN, "Fatigue < 0.40\nAND Relevance > 70%",
         "Both conditions must hold simultaneously across\nindependent runs.\n\n"
         "Fatigue = mean session_fatigue_accumulator\nat end of simulated sessions.\n\n"
         "Relevance = mean user satisfaction score\n(driven by ad-interest match).\n\n"
         "Reported as proportion of runs passing each\nthreshold independently + combined."),
        ("H3", ORANGE, "Diversity > 0.15",
         "The evolved policy uses a genuinely mixed\nstrategy across all 4 decision types.\n\n"
         "Measured as normalised Shannon entropy over\nSHOW / SOFTEN / DELAY / SUPPRESS.\n\n"
         "  D = H(decisions) / log₂(4)\n\n"
         "Threshold 0.15 rules out degenerate\nall-suppress or all-show policies."),
    ]

    for i, (label, color, title, body) in enumerate(hyps):
        l = Inches(0.35 + i * 4.3)
        card(s, l, Inches(1.55), Inches(4.1), Inches(5.7))
        chip(s, label, l + Inches(0.15), Inches(1.7), color=color, w=Inches(0.55))
        txt(s, title, l + Inches(0.8), Inches(1.68), Inches(3.1), Inches(0.55),
            size=14, bold=True, color=color)
        txt(s, body, l + Inches(0.15), Inches(2.35), Inches(3.82),
            Inches(4.75), size=11, color=SUB)
    return s


# ─────────────────────────────────────────────────────────────────────────────
# Slide 9 — Results
# ─────────────────────────────────────────────────────────────────────────────
def s09_results(prs):
    s = blank(prs); bg(s, prs)
    header(s, "RESULTS", "Quick test: 5 runs × 10 gen, 50 users  |  Full run: 30 × 50 gen pending", prs)

    # Baselines table
    card(s, Inches(0.35), Inches(1.55), Inches(6.5), Inches(3.15))
    txt(s, "Baseline comparison (200 users, 983 decisions)",
        Inches(0.55), Inches(1.65), Inches(6.1), Inches(0.4),
        size=13, bold=True, color=ACCENT)
    headers = ["Policy", "Satisfaction", "Revenue", "Fatigue", "Fitness"]
    rows = [
        ["always_show", "0.339", "0.770", "0.439", "0.512"],
        ["random",      "0.517", "0.408", "0.324", "0.473"],
        ["freq_cap_3",  "0.458", "0.467", "0.292", "0.462"],
        ["GA (quick)",  "0.—",   "0.—",   "0.330", "0.520"],
    ]
    col_x = [Inches(0.55), Inches(2.1), Inches(3.25), Inches(4.3), Inches(5.35)]
    row_colors = [WHITE, SUB, SUB, SUB, GREEN]
    for ri, row in enumerate([headers] + rows):
        for ci, cell in enumerate(row):
            top = Inches(2.1) + Inches(ri * 0.47)
            color = ACCENT if ri == 0 else (GREEN if ri == 4 and ci > 0 else SUB)
            txt(s, cell, col_x[ci], top, Inches(1.0), Inches(0.4),
                size=11, bold=(ri == 0), color=color)

    # H results
    h_results = [
        ("H1", RED,    "FAIL\n(quick test)",
         "0.52 mean < 0.58\nthreshold. GA beats\nall baselines but\nneeds 30×50 run\nto cross threshold."),
        ("H2", GREEN,  "PASS",
         "0.33 mean fatigue\n< 0.40 threshold.\nUser experience\nwell protected by\nsuppression policy."),
        ("H3", GREEN,  "PASS",
         "0.58 diversity\n> 0.15 threshold.\nHealthy mix of all\n4 decisions — not\na degenerate policy."),
    ]
    for i, (label, color, verdict, note) in enumerate(h_results):
        l = Inches(6.95 + i * 2.1)
        card(s, l, Inches(1.55), Inches(2.0), Inches(3.15))
        chip(s, label, l + Inches(0.1), Inches(1.68), color=color, w=Inches(0.55))
        txt(s, verdict, l + Inches(0.1), Inches(2.1), Inches(1.8), Inches(0.65),
            size=15, bold=True, color=color)
        txt(s, note, l + Inches(0.1), Inches(2.8), Inches(1.8), Inches(1.8),
            size=11, color=SUB)

    # GA trajectory
    card(s, Inches(0.35), Inches(4.9), Inches(12.6), Inches(1.0), fill=DEEP)
    txt(s, "GA trajectory  ·  quick test (10 generations):",
        Inches(0.55), Inches(4.98), Inches(4.0), Inches(0.38),
        size=12, bold=True, color=WHITE)
    txt(s, "Initial best: 0.501  →  Gen 10 best: 0.503  →  Full 30×50 expected to push past 0.58 threshold",
        Inches(0.55), Inches(5.4), Inches(12.1), Inches(0.38),
        size=12, color=SUB)

    card(s, Inches(0.35), Inches(6.05), Inches(12.6), Inches(1.15), fill=DEEP)
    txt(s, "A/B Human Evaluation (live sessions):",
        Inches(0.55), Inches(6.12), Inches(4.5), Inches(0.38),
        size=12, bold=True, color=WHITE)
    txt(s, "Score = Willingness + Relevance − Annoyance  (range −8 to +19)\n"
           "AdaptAd vs random baseline · sessions saved to SQLite · aggregate win/loss tracked across all participants",
        Inches(0.55), Inches(6.52), Inches(12.1), Inches(0.55),
        size=11, color=SUB)
    return s


# ─────────────────────────────────────────────────────────────────────────────
# Slide 10 — Hardcoded Decisions & Limitations
# ─────────────────────────────────────────────────────────────────────────────
def s10_limitations(prs):
    s = blank(prs); bg(s, prs)
    header(s, "LIMITATIONS & HARDCODED DECISIONS", "What we fixed, why we fixed it, and what it costs", prs)

    items = [
        (ORANGE, "Value Judgment: 60/40 fitness split",
         "fitness = 0.60 × satisfaction + 0.40 × revenue. "
         "No empirical basis — this is an ethical position. "
         "Different weights produce different evolved behaviours."),
        (ORANGE, "Value Judgment: 55/45 agent weights",
         "combined = 0.55 · UA + 0.45 · ADV. Separatefrom the fitness split. "
         "Neither weight is empirically derived — both are tunable but currently fixed in config."),
        (RED,    "Synthetic users — no real behavioural data",
         "1,000 users generated with age-weighted distributions. "
         "Real engagement patterns, fatigue curves, and tolerance "
         "distributions will differ. MovieLens grounds content genres; "
         "user psychology is modelled, not measured."),
        (RED,    "Fixed gene count = 8",
         "Chromosome length is hardcoded. The GA cannot discover "
         "that it needs more or fewer dimensions. Variable-length "
         "chromosomes would require a different evolutionary strategy."),
        (ORANGE, "Hardcoded mood & primetime modifiers",
         "MOOD_MODIFIER dict in fitness.py (calm=+0.10, dark=−0.15 etc.) "
         "and primetime map (evening=0.15) are estimates from literature, "
         "not learned from data. Wrong values silently skew the fitness landscape."),
        (ORANGE, "UserProfile.id = 99999 sentinel",
         "Custom A/B test profiles get id=99999 — a magic number "
         "that collides if two custom tests run simultaneously. "
         "No proper UUID generation for transient profiles."),
        (SUB,    "LLM fallback chain is brittle",
         "Groq → Gemini → template. If both API keys are missing "
         "or rate-limited, explanations silently fall back to "
         "deterministic templates with no indication in the UI."),
        (SUB,    "Fatigue increments are estimates",
         "show_increment=0.10, soften=0.05, delay=0.02 are hand-chosen. "
         "No user study backs these numbers. Wrong increments mean "
         "the force-suppress threshold triggers too early or too late."),
    ]

    for i, (color, title, body) in enumerate(items):
        col = 0 if i < 4 else 1
        row = i if i < 4 else i - 4
        l = Inches(0.35) if col == 0 else Inches(6.7)
        t = Inches(1.55) + row * Inches(1.37)
        card(s, l, t, Inches(6.15), Inches(1.28), fill=DEEP)
        sh = s.shapes.add_shape(1, l, t, Inches(0.07), Inches(1.28))
        sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
        txt(s, title, l + Inches(0.15), t + Inches(0.08), Inches(5.8),
            Inches(0.38), size=12, bold=True, color=color)
        txt(s, body, l + Inches(0.15), t + Inches(0.5), Inches(5.8),
            Inches(0.68), size=10, color=SUB)
    return s


# ─────────────────────────────────────────────────────────────────────────────
# Slide 11 — HCI & A/B Testing
# ─────────────────────────────────────────────────────────────────────────────
def s11_hci(prs):
    s = blank(prs); bg(s, prs)
    header(s, "HCI ANGLE", "Every decision explained. Humans in the loop.", prs)

    # Left — explainability
    card(s, Inches(0.35), Inches(1.55), Inches(6.15), Inches(5.65))
    chip(s, "Explainability", Inches(0.5), Inches(1.7), color=GREEN, w=Inches(2.2))
    explain_items = [
        "Every ad decision includes a natural-language reason.",
        "Example: \"Skipped — you've seen 3 ads this session,",
        "you're in an intense scene, and this ad isn't relevant",
        "to your interests.\"",
        "",
        "This is not a black box. Operators see exactly why",
        "the system chose each action.",
        "",
        "Explanation sources (in priority order):",
        "  1. Groq — Llama 3.3 70B (14,400 req/day free)",
        "  2. Gemini — 2.5 Flash (250 req/day free)",
        "  3. Template — deterministic offline fallback",
    ]
    txt(s, "\n".join(explain_items), Inches(0.5), Inches(2.2), Inches(5.7),
        Inches(4.8), size=11, color=SUB)

    # Right — A/B testing
    card(s, Inches(6.7), Inches(1.55), Inches(6.25), Inches(5.65))
    chip(s, "Human A/B Testing", Inches(6.85), Inches(1.7), color=ACCENT, w=Inches(2.5))

    steps = [
        ("Profile", "Participant fills name, age, occupation,\n"
                    "ad interests, genre prefs, ad tolerance,\nbinge tendency"),
        ("Content", "Enter show title — Auto-fill fetches genre,\n"
                    "duration, series/movie via LLM"),
        ("Blind test", "System generates Session X and Session Y —\n"
                       "participant does not know which is AdaptAd"),
        ("Rate", "Rate 1–10: Annoyance · Relevance · Would Continue"),
        ("Reveal", "Score = Willingness + Relevance − Annoyance\n"
                   "Winner revealed · saved to SQLite"),
    ]
    for i, (step, desc) in enumerate(steps):
        t = Inches(2.2) + i * Inches(0.95)
        sh = s.shapes.add_shape(1, Inches(6.85), t + Inches(0.08),
                                Inches(0.32), Inches(0.32))
        sh.fill.solid(); sh.fill.fore_color.rgb = ACCENT; sh.line.fill.background()
        tf = sh.text_frame; p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(i+1)
        r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = BG
        txt(s, step, Inches(7.28), t, Inches(1.5), Inches(0.4),
            size=11, bold=True, color=WHITE)
        txt(s, desc, Inches(7.28), t + Inches(0.38), Inches(5.4),
            Inches(0.55), size=10, color=SUB)
    return s


# ─────────────────────────────────────────────────────────────────────────────
# Slide 12 — Demo
# ─────────────────────────────────────────────────────────────────────────────
def s12_demo(prs):
    s = blank(prs); bg(s, prs)
    bar(s, prs, Inches(0.06)); vbar(s, prs)

    txt(s, "Live Demo", Inches(0.55), Inches(1.0), Inches(10), Inches(1.0),
        size=50, bold=True, color=ACCENT)
    txt(s, "uvicorn backend.main:app --port 8000  ·  npm run dev  →  localhost:5173",
        Inches(0.55), Inches(2.1), Inches(12), Inches(0.45), size=14, color=SUB)

    stops = [
        (Inches(0.55), Inches(2.8),  GREEN,  "① Evolution page",
         "Start a 10-generation run · watch fitness chart update live via WebSocket · "
         "see population grid and gene tracker animate · observe diversity metric"),
        (Inches(0.55), Inches(4.1),  ACCENT, "② Decision Explorer",
         "Pick a user + ad · run a single decision · read the agent scores and "
         "LLM explanation · change the user profile and see how the decision shifts"),
        (Inches(0.55), Inches(5.4),  ORANGE, "③ Session Simulator",
         "Simulate a full streaming session · see every ad break with its decision, "
         "combined score, and fatigue trajectory · compare evolved vs baseline policy"),
        (Inches(0.55), Inches(6.55), PINK,   "④ A/B Testing (live participant)",
         "Fill in your own profile · enter a real show · rate Session X vs Session Y "
         "blind · reveal which was AdaptAd · show win/loss in the saved history table"),
    ]
    for l, t, color, title, body in stops:
        sh = s.shapes.add_shape(1, l, t, Inches(0.06), Inches(1.0))
        sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
        txt(s, title, l + Inches(0.2), t, Inches(12.0), Inches(0.42),
            size=15, bold=True, color=color)
        txt(s, body, l + Inches(0.2), t + Inches(0.42), Inches(12.0),
            Inches(0.55), size=12, color=SUB)
    return s


# ─────────────────────────────────────────────────────────────────────────────
# Slide 13 — Conclusion
# ─────────────────────────────────────────────────────────────────────────────
def s13_conclusion(prs):
    s = blank(prs); bg(s, prs)
    bar(s, prs, Inches(0.06)); vbar(s, prs)

    txt(s, "AdaptAd", Inches(0.55), Inches(1.2), Inches(10), Inches(1.0),
        size=52, bold=True, color=ACCENT)
    txt(s, "Better ads.  Happier viewers.  Healthier platforms.",
        Inches(0.55), Inches(2.35), Inches(10), Inches(0.6),
        size=22, color=WHITE)

    takeaways = [
        (GREEN,  "What worked",
         "GA consistently outperformed all three baselines · "
         "H2 (fatigue) and H3 (diversity) both passed · "
         "LangGraph orchestration clean and extensible · "
         "WebSocket live evolution works end-to-end"),
        (ORANGE, "What's honest",
         "H1 needs the full 30×50 run to confirm · "
         "60/40 fitness split is a value judgment, not a fact · "
         "synthetic users are a model, not a measurement · "
         "LLM explanations are cosmetic if keys are missing"),
        (ACCENT, "What's next",
         "Run full experiment and write the paper · "
         "collect 5–10 A/B human participants · "
         "replace static fitness weights with a learned preference model · "
         "real user behaviour data from a streaming partner"),
    ]
    for i, (color, title, body) in enumerate(takeaways):
        l = Inches(0.35 + i * 4.3)
        card(s, l, Inches(3.3), Inches(4.1), Inches(3.6))
        chip(s, title, l + Inches(0.15), Inches(3.45), color=color, w=Inches(2.0))
        txt(s, body, l + Inches(0.15), Inches(3.95), Inches(3.82),
            Inches(2.8), size=11, color=SUB)

    txt(s, "github · uvicorn backend.main:app --port 8000 · cd frontend && npm run dev",
        Inches(0.55), Inches(7.1), Inches(10), Inches(0.32),
        size=11, color=RGBColor(0x2a,0x40,0x58), italic=True)

    GCOLORS = [ACCENT, GREEN, ORANGE, RED, PURPLE, PINK,
               RGBColor(0x34,0xd3,0x99), RGBColor(0xfb,0x92,0x3c)]
    vals = [0.46, 0.92, 0.35, 0.39, 0.08, 0.80, 0.00, 0.66]
    for i, (c, v) in enumerate(zip(GCOLORS, vals)):
        w = Inches(0.08 + 1.9 * v)
        sh = s.shapes.add_shape(1, Inches(10.9), Inches(1.2 + i*0.5), w, Inches(0.3))
        sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
    return s


# ─────────────────────────────────────────────────────────────────────────────
# Build
# ─────────────────────────────────────────────────────────────────────────────
def build():
    prs = prs_new()
    s01_title(prs)
    s02_problem(prs)
    s03_answer(prs)
    s04_arch(prs)
    s05_chrom(prs)
    s06_agents(prs)
    s07_ga(prs)
    s08_hypo(prs)
    s09_results(prs)
    s10_limitations(prs)
    s11_hci(prs)
    s12_demo(prs)
    s13_conclusion(prs)

    out = "AdaptAd_Presentation_v2.pptx"
    prs.save(out)
    print(f"✓ Saved {len(prs.slides)} slides → {out}")

if __name__ == "__main__":
    build()
