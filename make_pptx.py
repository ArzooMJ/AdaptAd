"""
Generate AdaptAd presentation as PPTX.
Run: python3 make_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Palette ───────────────────────────────────────────────────────────────────
BG        = RGBColor(0x0d, 0x15, 0x1f)   # dark navy
ACCENT    = RGBColor(0x00, 0xd4, 0xff)   # cyan
GREEN     = RGBColor(0x00, 0xff, 0x88)   # mint green
ORANGE    = RGBColor(0xff, 0xb8, 0x00)   # amber
RED       = RGBColor(0xff, 0x2d, 0x55)   # red
WHITE     = RGBColor(0xff, 0xff, 0xff)
SUBTEXT   = RGBColor(0x8b, 0x9a, 0xb2)   # muted blue-grey
CARD      = RGBColor(0x14, 0x21, 0x33)   # card background


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def bg_rect(slide, prs):
    """Full-slide dark background."""
    w, h = prs.slide_width, prs.slide_height
    shape = slide.shapes.add_shape(1, 0, 0, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def accent_bar(slide, prs, height=Inches(0.05)):
    """Thin cyan bar along the top."""
    w = prs.slide_width
    shape = slide.shapes.add_shape(1, 0, 0, w, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def section_chip(slide, text, left, top, color=ACCENT):
    """Small coloured label chip."""
    w, h = Inches(1.6), Inches(0.32)
    shape = slide.shapes.add_shape(1, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = BG


def card_rect(slide, left, top, width, height, color=CARD):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = RGBColor(0x20, 0x35, 0x50)
    shape.line.width = Pt(0.75)
    return shape


# ─────────────────────────────────────────────────────────────────────────────
# Slides
# ─────────────────────────────────────────────────────────────────────────────

def slide_title(prs):
    s = blank_slide(prs)
    bg_rect(s, prs)
    accent_bar(s, prs, Inches(0.06))

    # Vertical cyan stripe
    stripe = s.shapes.add_shape(1, 0, 0, Inches(0.06), prs.slide_height)
    stripe.fill.solid(); stripe.fill.fore_color.rgb = ACCENT
    stripe.line.fill.background()

    add_text(s, "AdaptAd", Inches(0.5), Inches(1.6), Inches(9), Inches(1.2),
             size=54, bold=True, color=ACCENT)
    add_text(s, "Human-Centered Ad Intelligence for Streaming Platforms",
             Inches(0.5), Inches(2.9), Inches(9.5), Inches(0.7),
             size=22, color=WHITE)
    add_text(s, "CS6170 AI Capstone  ·  Northeastern University",
             Inches(0.5), Inches(3.7), Inches(9), Inches(0.5),
             size=14, color=SUBTEXT)
    add_text(s, "Craig Roberts  ·  Arzoo Jiwani  ·  Vishwajeet Hogale",
             Inches(0.5), Inches(4.15), Inches(9), Inches(0.45),
             size=13, color=SUBTEXT)

    # Decorative gene bars on the right
    colors = [ACCENT, GREEN, ORANGE, RED,
              RGBColor(0xa7,0x8b,0xfa), RGBColor(0xf4,0x72,0xb6),
              RGBColor(0x34,0xd3,0x99), RGBColor(0xfb,0x92,0x3c)]
    values = [0.46, 0.92, 0.35, 0.39, 0.08, 0.80, 0.00, 0.66]
    bar_left = Inches(10.5)
    bar_top  = Inches(2.0)
    for i, (c, v) in enumerate(zip(colors, values)):
        bar_w = Inches(1.8 * v + 0.1)
        shape = s.shapes.add_shape(1, bar_left, bar_top + Inches(i * 0.42),
                                   bar_w, Inches(0.28))
        shape.fill.solid(); shape.fill.fore_color.rgb = c
        shape.line.fill.background()

    add_text(s, "Best chromosome · 8 genes", Inches(10.4), Inches(5.6),
             Inches(2.5), Inches(0.4), size=10, color=SUBTEXT, italic=True)
    return s


def slide_problem(prs):
    s = blank_slide(prs)
    bg_rect(s, prs); accent_bar(s, prs)
    section_chip(s, "THE PROBLEM", Inches(0.5), Inches(0.25))
    add_text(s, "Streaming ads are broken", Inches(0.5), Inches(0.75),
             Inches(9), Inches(0.8), size=34, bold=True, color=WHITE)

    problems = [
        ("Ad fatigue is real",
         "72% of viewers say they feel overwhelmed by streaming ads.\nRepetitive, ill-timed ads drive subscription cancellations."),
        ("Click-through rate is the wrong goal",
         "Optimising for CTR ignores whether the viewer was in the right\nstate to receive an ad — and punishes long-term retention."),
        ("One-size-fits-all scheduling",
         "Static frequency caps treat a binge-watcher finishing a\ncliffhanger the same as a casual viewer in a calm scene."),
    ]

    for i, (title, body) in enumerate(problems):
        left = Inches(0.4 + i * 4.25)
        card_rect(s, left, Inches(1.8), Inches(4.0), Inches(4.5))
        add_text(s, title, left + Inches(0.2), Inches(2.0), Inches(3.6),
                 Inches(0.6), size=15, bold=True, color=ACCENT)
        add_text(s, body, left + Inches(0.2), Inches(2.7), Inches(3.6),
                 Inches(3.2), size=12, color=SUBTEXT)
    return s


def slide_solution(prs):
    s = blank_slide(prs)
    bg_rect(s, prs); accent_bar(s, prs)
    section_chip(s, "OUR APPROACH", Inches(0.5), Inches(0.25))
    add_text(s, "Ask 'should we even show an ad?' not 'which ad gets more clicks?'",
             Inches(0.5), Inches(0.75), Inches(12.3), Inches(0.8),
             size=26, bold=True, color=WHITE)

    decisions = [
        ("SHOW",     GREEN,  "Conditions are favorable.\nFull ad."),
        ("SOFTEN",   ACCENT, "Moderate fit.\nShorter version."),
        ("DELAY",    ORANGE, "Bad timing, good ad.\nWait for better moment."),
        ("SUPPRESS", RED,    "Protect the viewer.\nSkip entirely."),
    ]

    for i, (label, color, desc) in enumerate(decisions):
        left = Inches(0.4 + i * 3.2)
        card_rect(s, left, Inches(1.8), Inches(3.0), Inches(4.2))
        chip = s.shapes.add_shape(1, left + Inches(0.15), Inches(2.0),
                                  Inches(1.6), Inches(0.38))
        chip.fill.solid(); chip.fill.fore_color.rgb = color
        chip.line.fill.background()
        tf = chip.text_frame; p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = label
        r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = BG
        add_text(s, desc, left + Inches(0.15), Inches(2.55), Inches(2.7),
                 Inches(2.8), size=12, color=SUBTEXT)

    add_text(s,
             "For every ad opportunity in a streaming session, AdaptAd picks one of the four actions above "
             "using a Genetic Algorithm–evolved chromosome and a two-agent scoring system.",
             Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.9),
             size=13, color=SUBTEXT)
    return s


def slide_architecture(prs):
    s = blank_slide(prs)
    bg_rect(s, prs); accent_bar(s, prs)
    section_chip(s, "ARCHITECTURE", Inches(0.5), Inches(0.25))
    add_text(s, "System overview", Inches(0.5), Inches(0.75),
             Inches(8), Inches(0.6), size=30, bold=True, color=WHITE)

    # Flow diagram as styled boxes + arrows
    boxes = [
        (Inches(0.35), "Streaming\nSession",   CARD,  ACCENT),
        (Inches(2.45), "Ad\nOpportunity",       CARD,  ACCENT),
        (Inches(4.55), "User Advocate\n(math)", CARD,  GREEN),
        (Inches(6.65), "Advertiser\nAdvocate",  CARD,  ORANGE),
        (Inches(8.75), "Negotiator\n(GA genes)",CARD,  ACCENT),
        (Inches(10.85),"Decision\nSHOW·…·SUPP",CARD,  RED),
    ]
    for left, label, bg, outline in boxes:
        card = s.shapes.add_shape(1, left, Inches(2.0), Inches(1.9), Inches(1.3))
        card.fill.solid(); card.fill.fore_color.rgb = bg
        card.line.color.rgb = outline; card.line.width = Pt(1.5)
        tf = card.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = label
        r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = WHITE

    # Arrows
    for x in [Inches(2.26), Inches(4.36), Inches(6.46), Inches(8.56), Inches(10.66)]:
        arr = s.shapes.add_shape(1, x, Inches(2.5), Inches(0.19), Inches(0.3))
        arr.fill.solid(); arr.fill.fore_color.rgb = SUBTEXT
        arr.line.fill.background()

    # GA box below
    card_rect(s, Inches(3.5), Inches(3.8), Inches(6.3), Inches(2.1),
              RGBColor(0x08, 0x14, 0x28))
    add_text(s, "Genetic Algorithm Evolution", Inches(3.7), Inches(3.95),
             Inches(5.9), Inches(0.5), size=14, bold=True, color=ACCENT)
    add_text(s,
             "Population of 30 chromosomes  ·  Fitness = 60% satisfaction + 40% revenue\n"
             "Tournament selection  ·  Uniform crossover  ·  Gaussian mutation  ·  Elite preservation\n"
             "Convergence detection  ·  Warm restart when stuck  ·  Best chromosome saved to disk",
             Inches(3.7), Inches(4.5), Inches(5.9), Inches(1.3),
             size=11, color=SUBTEXT)

    # LLM box
    card_rect(s, Inches(0.35), Inches(3.8), Inches(2.85), Inches(2.1),
              RGBColor(0x08, 0x14, 0x28))
    add_text(s, "LLM Explain", Inches(0.5), Inches(3.95), Inches(2.6),
             Inches(0.5), size=14, bold=True, color=RGBColor(0xa7,0x8b,0xfa))
    add_text(s,
             "Groq (primary)\nGemini (fallback)\nTemplate (offline)\n\nNatural language\ndecision reasoning",
             Inches(0.5), Inches(4.5), Inches(2.6), Inches(1.3),
             size=11, color=SUBTEXT)
    return s


def slide_chromosome(prs):
    s = blank_slide(prs)
    bg_rect(s, prs); accent_bar(s, prs)
    section_chip(s, "8-GENE CHROMOSOME", Inches(0.5), Inches(0.25))
    add_text(s, "All 8 genes are active in the fitness landscape",
             Inches(0.5), Inches(0.75), Inches(9), Inches(0.65),
             size=28, bold=True, color=WHITE)

    genes = [
        ("fatigue_weight",       "0.46", "Sensitivity to session fatigue",    "More conservative for tired users"),
        ("relevance_weight",     "0.92", "Importance of ad-interest match",    "Only shows ads to users whose interests align"),
        ("timing_weight",        "0.35", "Time-of-day alignment importance",   "Favours preferred viewing times"),
        ("frequency_threshold",  "0.39", "Base bar to show any ad",            "Show threshold maps to range 0.35–0.65"),
        ("delay_probability",    "0.08", "Width of the DELAY zone",            "Prefers delaying over suppressing"),
        ("soften_threshold",     "0.80", "Width of the SOFTEN zone",           "Shorter ads over full skip when borderline"),
        ("category_boost",       "0.00", "Advertiser category relevance wt.",  "Rewards relevant ads more for advertisers"),
        ("session_depth_factor", "0.66", "Penalty growth with ads_shown",      "Increasingly cautious deep in a session"),
    ]

    colors = [ACCENT, GREEN, ORANGE, RED,
              RGBColor(0xa7,0x8b,0xfa), RGBColor(0xf4,0x72,0xb6),
              RGBColor(0x34,0xd3,0x99), RGBColor(0xfb,0x92,0x3c)]

    row_h = Inches(0.55)
    for i, (name, val, short, long_) in enumerate(genes):
        top = Inches(1.65) + i * row_h
        color = colors[i]

        # Gene bar
        bar_w = Inches(float(val) * 1.5 + 0.05)
        bar = s.shapes.add_shape(1, Inches(0.4), top + Inches(0.12),
                                 bar_w, Inches(0.25))
        bar.fill.solid(); bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        add_text(s, val, Inches(0.4), top, Inches(0.6), Inches(0.5),
                 size=10, bold=True, color=color)
        add_text(s, name, Inches(2.1), top, Inches(3.2), Inches(0.5),
                 size=11, bold=True, color=WHITE)
        add_text(s, short, Inches(5.4), top, Inches(3.5), Inches(0.5),
                 size=10, color=SUBTEXT)
        add_text(s, long_, Inches(9.0), top, Inches(4.0), Inches(0.5),
                 size=10, color=RGBColor(0x55, 0x6a, 0x82), italic=True)

    add_text(s,
             "UserProfile.ad_tolerance is also used: high-tolerance users get higher satisfaction "
             "under the same SHOW decision, creating realistic heterogeneity across 200 users.",
             Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.4),
             size=11, color=SUBTEXT, italic=True)
    return s


def slide_agents(prs):
    s = blank_slide(prs)
    bg_rect(s, prs); accent_bar(s, prs)
    section_chip(s, "AGENT SYSTEM", Inches(0.5), Inches(0.25))
    add_text(s, "Two advocates, one negotiator",
             Inches(0.5), Inches(0.75), Inches(9), Inches(0.65),
             size=30, bold=True, color=WHITE)

    panels = [
        ("User Advocate", GREEN,
         ["Scores how receptive the viewer is right now",
          "Inputs: fatigue, ad relevance, time-of-day match,",
          "  content mood, binge state, session depth",
          "Gene weights: fatigue_weight, relevance_weight,",
          "  timing_weight, session_depth_factor",
          "Output: UA score ∈ [0,1]"]),
        ("Advertiser Advocate", ORANGE,
         ["Scores advertiser value of showing an ad here",
          "Inputs: category match, user engagement,",
          "  primetime slot, ad priority, seasonal affinity,",
          "  demographic alignment",
          "Gene weights: category_boost",
          "Output: ADV score ∈ [0,1]"]),
        ("Negotiator", ACCENT,
         ["Combines UA + ADV weighted by config",
          "combined = 0.6·UA + 0.4·ADV",
          "Maps combined score to a decision via",
          "  frequency_threshold → show_thresh",
          "  soften_threshold → soften zone width",
          "  delay_probability → delay zone width",
          "Output: SHOW / SOFTEN / DELAY / SUPPRESS"]),
    ]

    for i, (title, color, lines) in enumerate(panels):
        left = Inches(0.35 + i * 4.3)
        card_rect(s, left, Inches(1.7), Inches(4.1), Inches(5.3))
        chip = s.shapes.add_shape(1, left + Inches(0.15), Inches(1.85),
                                  Inches(2.4), Inches(0.36))
        chip.fill.solid(); chip.fill.fore_color.rgb = color
        chip.line.fill.background()
        tf = chip.text_frame; p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = title
        r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = BG

        body = "\n".join(lines)
        add_text(s, body, left + Inches(0.18), Inches(2.3), Inches(3.75),
                 Inches(4.5), size=11, color=SUBTEXT)
    return s


def slide_ga(prs):
    s = blank_slide(prs)
    bg_rect(s, prs); accent_bar(s, prs)
    section_chip(s, "GENETIC ALGORITHM", Inches(0.5), Inches(0.25))
    add_text(s, "Evolution pipeline", Inches(0.5), Inches(0.75),
             Inches(8), Inches(0.6), size=30, bold=True, color=WHITE)

    steps = [
        ("Init",       ACCENT,  "30 random\nchromosomes\nUniform [0,1]"),
        ("Evaluate",   GREEN,   "Vectorised NumPy\nfitness on 200 users\n5 scenarios each"),
        ("Select",     ORANGE,  "3-way tournament\nselection for\nparent pairs"),
        ("Crossover",  ACCENT,  "Uniform crossover:\neach gene from\nA or B at 50%"),
        ("Mutate",     RED,     "Gaussian Δ per gene\nat mutation_rate\nclamped [0,1]"),
        ("Elite",      GREEN,   "Top 20% survive\nunchanged into\nnext generation"),
        ("Converge",   SUBTEXT, "Δbest < 0.001 for\n15 generations\nor max gen reached"),
    ]

    box_w = Inches(1.7)
    for i, (label, color, desc) in enumerate(steps):
        left = Inches(0.3 + i * 1.85)
        card_rect(s, left, Inches(1.8), box_w, Inches(3.5))
        chip = s.shapes.add_shape(1, left + Inches(0.1), Inches(1.95),
                                  Inches(1.5), Inches(0.34))
        chip.fill.solid(); chip.fill.fore_color.rgb = color
        chip.line.fill.background()
        tf = chip.text_frame; p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = label
        r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = BG

        add_text(s, desc, left + Inches(0.1), Inches(2.4), Inches(1.5),
                 Inches(2.8), size=10, color=SUBTEXT, align=PP_ALIGN.CENTER)

    # Arrows between boxes
    for x_idx in range(6):
        ax = Inches(1.9 + x_idx * 1.85)
        arr = s.shapes.add_shape(1, ax, Inches(2.9), Inches(0.18), Inches(0.22))
        arr.fill.solid(); arr.fill.fore_color.rgb = SUBTEXT
        arr.line.fill.background()

    # Loop-back arrow text
    add_text(s, "↩ repeat until convergence", Inches(0.3), Inches(5.6),
             Inches(12.5), Inches(0.4), size=12, color=SUBTEXT, italic=True,
             align=PP_ALIGN.CENTER)

    # Config panel
    card_rect(s, Inches(0.3), Inches(6.1), Inches(12.7), Inches(1.0))
    add_text(s,
             "Config: population=30  ·  elite_ratio=0.20  ·  mutation_rate=0.20  ·  "
             "mutation_strength=0.15  ·  convergence_window=15  ·  convergence_threshold=0.001  ·  "
             "stuck_restart_threshold=20  ·  fitness = 0.60·satisfaction + 0.40·revenue",
             Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.8),
             size=11, color=SUBTEXT)
    return s


def slide_hypotheses(prs):
    s = blank_slide(prs)
    bg_rect(s, prs); accent_bar(s, prs)
    section_chip(s, "HYPOTHESES", Inches(0.5), Inches(0.25))
    add_text(s, "Three testable predictions", Inches(0.5), Inches(0.75),
             Inches(9), Inches(0.6), size=30, bold=True, color=WHITE)

    hyps = [
        ("H1", ACCENT,  "Evolved fitness > 0.58",
         "The GA-evolved policy achieves a mean composite fitness score significantly above 0.58 "
         "and outperforms all three baselines (always-show, random, freq-cap-3).\n\n"
         "Test: one-sample Wilcoxon signed-rank + paired Wilcoxon vs each baseline, "
         "Holm–Bonferroni corrected."),
        ("H2", GREEN,   "Fatigue < 0.40  AND  Relevance > 70%",
         "Post-session ad fatigue stays below 0.40 AND mean user satisfaction "
         "(proxy for relevance) exceeds 0.70.\n\n"
         "Both conditions must hold simultaneously. Reported as proportion of runs passing "
         "each threshold."),
        ("H3", ORANGE,  "Strategy diversity > 0.15",
         "The evolved policy uses a genuinely diverse mix of SHOW/SOFTEN/DELAY/SUPPRESS decisions, "
         "measured as normalised Shannon entropy over the 4 decision types.\n\n"
         "Threshold 0.15 on a [0,1] scale — eliminates degenerate all-suppress / all-show policies."),
    ]

    for i, (label, color, title, body) in enumerate(hyps):
        left = Inches(0.35 + i * 4.3)
        card_rect(s, left, Inches(1.75), Inches(4.1), Inches(5.35))
        chip = s.shapes.add_shape(1, left + Inches(0.15), Inches(1.9),
                                  Inches(0.5), Inches(0.36))
        chip.fill.solid(); chip.fill.fore_color.rgb = color
        chip.line.fill.background()
        tf = chip.text_frame; p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = label
        r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = BG

        add_text(s, title, left + Inches(0.75), Inches(1.88), Inches(3.25),
                 Inches(0.45), size=13, bold=True, color=color)
        add_text(s, body, left + Inches(0.15), Inches(2.45), Inches(3.8),
                 Inches(4.5), size=11, color=SUBTEXT)
    return s


def slide_results(prs):
    s = blank_slide(prs)
    bg_rect(s, prs); accent_bar(s, prs)
    section_chip(s, "PRELIMINARY RESULTS", Inches(0.5), Inches(0.25))
    add_text(s, "5 runs × 10 generations, 50 users  (quick test — full run = 30×50)",
             Inches(0.5), Inches(0.75), Inches(12.3), Inches(0.6),
             size=22, bold=True, color=WHITE)

    # Baselines table
    card_rect(s, Inches(0.35), Inches(1.6), Inches(5.8), Inches(2.9))
    add_text(s, "Baselines (200 users, 983 decisions)",
             Inches(0.55), Inches(1.7), Inches(5.4), Inches(0.45),
             size=13, bold=True, color=ACCENT)
    rows = [
        ("Policy",       "Satisfaction", "Revenue", "Fatigue", "Fitness"),
        ("always_show",  "0.339",        "0.770",   "0.439",   "0.512"),
        ("random",       "0.517",        "0.408",   "0.324",   "0.473"),
        ("freq_cap_3",   "0.458",        "0.467",   "0.292",   "0.462"),
    ]
    col_x = [Inches(0.55), Inches(2.05), Inches(3.1), Inches(4.1), Inches(5.05)]
    for r_i, row in enumerate(rows):
        for c_i, cell in enumerate(row):
            top = Inches(2.2) + Inches(0.45) * r_i
            is_hdr = r_i == 0
            add_text(s, cell, col_x[c_i], top, Inches(0.9), Inches(0.4),
                     size=11, bold=is_hdr,
                     color=SUBTEXT if is_hdr else WHITE)

    # H1/H2/H3 results
    results = [
        ("H1 — Fitness", RED,    "FAIL  0.52 mean",
         "Threshold: 0.58\nNeeds 30 runs × 50 gen for paper."),
        ("H2 — Fatigue", GREEN,  "PASS  0.33 mean",
         "Threshold: < 0.40\nFatigue well controlled."),
        ("H3 — Diversity",GREEN, "PASS  0.58 mean",
         "Threshold: > 0.15\nHealthy strategy mix."),
    ]
    for i, (title, color, verdict, note) in enumerate(results):
        left = Inches(6.4 + i * 2.25)
        card_rect(s, left, Inches(1.6), Inches(2.1), Inches(2.9))
        add_text(s, title, left + Inches(0.12), Inches(1.75), Inches(1.85),
                 Inches(0.5), size=11, bold=True, color=color)
        add_text(s, verdict, left + Inches(0.12), Inches(2.3), Inches(1.85),
                 Inches(0.55), size=14, bold=True, color=color)
        add_text(s, note, left + Inches(0.12), Inches(2.9), Inches(1.85),
                 Inches(1.4), size=10, color=SUBTEXT)

    # GA improvement note
    card_rect(s, Inches(0.35), Inches(4.7), Inches(12.6), Inches(1.4))
    add_text(s, "GA trajectory (10 generations)",
             Inches(0.55), Inches(4.82), Inches(5), Inches(0.4),
             size=12, bold=True, color=ACCENT)
    add_text(s,
             "Initial best fitness: 0.501   →   Gen 10 best: 0.503\n"
             "Early improvement is slow — the GA needs 30–50 generations to fully exploit the landscape. "
             "Full paper run (30×50) expected to push fitness above 0.58 threshold.",
             Inches(0.55), Inches(5.25), Inches(12.1), Inches(0.7),
             size=11, color=SUBTEXT)
    return s


def slide_ab(prs):
    s = blank_slide(prs)
    bg_rect(s, prs); accent_bar(s, prs)
    section_chip(s, "A/B TESTING", Inches(0.5), Inches(0.25))
    add_text(s, "Human evaluation — blind comparison",
             Inches(0.5), Inches(0.75), Inches(9), Inches(0.6),
             size=30, bold=True, color=WHITE)

    steps = [
        ("1  Profile",   "Participant fills in:\nage group · occupation\nad interests · genre prefs\nad tolerance · binge tendency"),
        ("2  Content",   "Participant enters\nshow/movie title.\nAuto-fill fetches genre,\nduration, series/movie."),
        ("3  Two schedules", "System generates:\nSession X — one policy\nSession Y — other policy\nLabel assignment is random."),
        ("4  Blind rating", "Rate each session 1–10 on:\nAnnoyance\nRelevance\nWould Continue?"),
        ("5  Reveal",    "Score = Willingness\n+ Relevance − Annoyance\nWinner revealed\nall sessions saved to DB."),
    ]
    for i, (title, body) in enumerate(steps):
        left = Inches(0.35 + i * 2.55)
        card_rect(s, left, Inches(1.7), Inches(2.4), Inches(4.0))
        add_text(s, title, left + Inches(0.15), Inches(1.85), Inches(2.1),
                 Inches(0.5), size=12, bold=True, color=ACCENT)
        add_text(s, body, left + Inches(0.15), Inches(2.45), Inches(2.1),
                 Inches(3.1), size=11, color=SUBTEXT)

    add_text(s,
             "AdaptAd vs random-baseline  ·  neither participant nor experimenter knows which session "
             "is AI-optimised during rating  ·  aggregate win/loss tracked across all participants",
             Inches(0.35), Inches(6.0), Inches(12.6), Inches(0.7),
             size=12, color=SUBTEXT, italic=True)
    return s


def slide_tech_stack(prs):
    s = blank_slide(prs)
    bg_rect(s, prs); accent_bar(s, prs)
    section_chip(s, "TECH STACK", Inches(0.5), Inches(0.25))
    add_text(s, "Full-stack AI system, built from scratch",
             Inches(0.5), Inches(0.75), Inches(9), Inches(0.6),
             size=30, bold=True, color=WHITE)

    cols = [
        ("Backend", ACCENT, [
            "Python 3.10+",
            "FastAPI + uvicorn",
            "Pydantic v2 models",
            "NumPy vectorised GA",
            "LangGraph pipelines",
            "SQLite + aiosqlite",
            "WebSocket real-time",
            "Groq / Gemini LLM",
        ]),
        ("Frontend", GREEN, [
            "React 18 + TypeScript",
            "Vite build tool",
            "Tailwind CSS",
            "Zustand state",
            "Recharts charts",
            "WebSocket hook",
            "Dark terminal theme",
            "Light mode support",
        ]),
        ("Data & Eval", ORANGE, [
            "200 synthetic users",
            "80 ads · 8 categories",
            "100 content items",
            "MovieLens 25M (optional)",
            "Criteo Display Ads",
            "Avazu CTR dataset",
            "Wilcoxon signed-rank",
            "Holm–Bonferroni corr.",
        ]),
        ("Scale", SUBTEXT, [
            "67 source files",
            "28 / 28 tests passing",
            "8 API route modules",
            "5 React pages",
            "~8 min per GA run",
            "Deployable to Render",
            "< 200ms API latency",
            "No GPU required",
        ]),
    ]
    for i, (title, color, items) in enumerate(cols):
        left = Inches(0.35 + i * 3.2)
        card_rect(s, left, Inches(1.65), Inches(3.05), Inches(5.55))
        chip = s.shapes.add_shape(1, left + Inches(0.15), Inches(1.8),
                                  Inches(2.0), Inches(0.32))
        chip.fill.solid(); chip.fill.fore_color.rgb = color
        chip.line.fill.background()
        tf = chip.text_frame; p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = title
        r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = BG

        body = "\n".join(f"• {x}" for x in items)
        add_text(s, body, left + Inches(0.15), Inches(2.2), Inches(2.75),
                 Inches(4.85), size=11, color=SUBTEXT)
    return s


def slide_future(prs):
    s = blank_slide(prs)
    bg_rect(s, prs); accent_bar(s, prs)
    section_chip(s, "NEXT STEPS", Inches(0.5), Inches(0.25))
    add_text(s, "Road to the paper deadline",
             Inches(0.5), Inches(0.75), Inches(9), Inches(0.6),
             size=30, bold=True, color=WHITE)

    items = [
        (GREEN,  "Run full 30×50 experiment",
                 "30 independent GA runs × 50 generations on all 200 users. "
                 "Expected runtime ~4 hours. Results saved to results/experiment_{timestamp}.json."),
        (ACCENT, "Collect 5–10 A/B participants",
                 "Human evaluation with real participants filling the custom profile form. "
                 "Aggregate win/loss statistics inform H1 from a human-preference angle."),
        (ORANGE, "Statistical analysis",
                 "Wilcoxon signed-rank tests (scipy) with Holm–Bonferroni correction. "
                 "Report p-values for H1 one-sample test and pairwise vs each baseline."),
        (RED,    "Write the paper",
                 "4–6 page CHI/CSCW-style paper covering system design, experiment, "
                 "results, and discussion. Deadline April 13 2026."),
    ]

    for i, (color, title, body) in enumerate(items):
        left  = Inches(0.35) if i < 2 else Inches(6.6)
        top   = Inches(1.75) if i % 2 == 0 else Inches(4.3)
        card_rect(s, left, top, Inches(6.0), Inches(2.3))
        add_text(s, title, left + Inches(0.2), top + Inches(0.15), Inches(5.6),
                 Inches(0.5), size=15, bold=True, color=color)
        add_text(s, body, left + Inches(0.2), top + Inches(0.7), Inches(5.6),
                 Inches(1.45), size=12, color=SUBTEXT)
    return s


def slide_closing(prs):
    s = blank_slide(prs)
    bg_rect(s, prs)
    accent_bar(s, prs, Inches(0.06))

    stripe = s.shapes.add_shape(1, 0, 0, Inches(0.06), prs.slide_height)
    stripe.fill.solid(); stripe.fill.fore_color.rgb = ACCENT
    stripe.line.fill.background()

    add_text(s, "AdaptAd", Inches(0.5), Inches(2.0), Inches(9), Inches(1.1),
             size=52, bold=True, color=ACCENT)
    add_text(s, "Better ads. Happier viewers. Healthier platforms.",
             Inches(0.5), Inches(3.2), Inches(9), Inches(0.7),
             size=22, color=WHITE)
    add_text(s,
             "Evolving ad policy chromosomes that respect viewer state —\n"
             "not just impressions and clicks.",
             Inches(0.5), Inches(4.1), Inches(9), Inches(0.9),
             size=15, color=SUBTEXT)
    add_text(s, "github · uvicorn backend.main:app --port 8000 · npm run dev",
             Inches(0.5), Inches(5.2), Inches(9), Inches(0.5),
             size=12, color=RGBColor(0x2a, 0x40, 0x58), italic=True)

    # Gene bars
    colors = [ACCENT, GREEN, ORANGE, RED,
              RGBColor(0xa7,0x8b,0xfa), RGBColor(0xf4,0x72,0xb6),
              RGBColor(0x34,0xd3,0x99), RGBColor(0xfb,0x92,0x3c)]
    values = [0.46, 0.92, 0.35, 0.39, 0.08, 0.80, 0.00, 0.66]
    for i, (c, v) in enumerate(zip(colors, values)):
        bar_w = Inches(1.8 * v + 0.1)
        shape = s.shapes.add_shape(1, Inches(10.5),
                                   Inches(2.0) + Inches(i * 0.42),
                                   bar_w, Inches(0.28))
        shape.fill.solid(); shape.fill.fore_color.rgb = c
        shape.line.fill.background()
    return s


# ─────────────────────────────────────────────────────────────────────────────
# Build
# ─────────────────────────────────────────────────────────────────────────────

def build():
    prs = new_prs()
    slide_title(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_architecture(prs)
    slide_chromosome(prs)
    slide_agents(prs)
    slide_ga(prs)
    slide_hypotheses(prs)
    slide_results(prs)
    slide_ab(prs)
    slide_tech_stack(prs)
    slide_future(prs)
    slide_closing(prs)

    out = "AdaptAd_Presentation.pptx"
    prs.save(out)
    print(f"Saved {len(prs.slides)} slides → {out}")


if __name__ == "__main__":
    build()
